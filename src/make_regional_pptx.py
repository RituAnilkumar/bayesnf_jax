"""
src/make_regional_pptx.py

Build outputs/egu_fin/regional_no_te.pptx
  1 title slide + 19 regions × 3 slides (A: time series, B: bar chart, C: beeswarm)

Usage:
    python src/make_regional_pptx.py \
        --ensemble_root outputs/egu_fin/ensemble_te_r2 \
        --explain_root  outputs/egu_fin/explain_te \
        --output        outputs/egu_fin/regional_no_te.pptx
"""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

# ---------------------------------------------------------------------------
# Layout constants (widescreen 13.33" × 7.5")
# ---------------------------------------------------------------------------
SW = Inches(13.333)
SH = Inches(7.5)

# Title bar
TTL_L, TTL_T, TTL_W, TTL_H = Inches(0), Inches(0), SW, Inches(1.0)
# Image (left panel)
IMG_L, IMG_T, IMG_W, IMG_H = Inches(0.25), Inches(1.1), Inches(8.4), Inches(6.1)
# Bullet box (right panel)
BUL_L, BUL_T, BUL_W, BUL_H = Inches(8.8), Inches(1.1), Inches(4.3), Inches(6.1)
# Footer strip
FTR_L, FTR_T, FTR_W, FTR_H = Inches(0), Inches(7.1), SW, Inches(0.4)

# ---------------------------------------------------------------------------
# Colors
# ---------------------------------------------------------------------------
NAVY    = RGBColor(0x1a, 0x23, 0x4b)
WHITE   = RGBColor(0xff, 0xff, 0xff)
TEAL    = RGBColor(0x2e, 0x86, 0xab)
OFFWHITE= RGBColor(0xf7, 0xf7, 0xf9)
DARK    = RGBColor(0x22, 0x22, 0x33)
SILVER  = RGBColor(0xb0, 0xb8, 0xc8)
GROUP   = RGBColor(0x8e, 0xc8, 0xda)   # for group labels in bullets

# ---------------------------------------------------------------------------
# Region metadata
# ---------------------------------------------------------------------------
REGIONS = [
    ("r01", "Alaska"),
    ("r02", "Western Canada & US"),
    ("r03", "Arctic Canada North"),
    ("r04", "Arctic Canada South"),
    ("r05", "Greenland Periphery"),
    ("r06", "Iceland"),
    ("r07", "Svalbard"),
    ("r08", "Scandinavia"),
    ("r09", "Russian Arctic"),
    ("r10", "North Asia"),
    ("r11", "Central Europe"),
    ("r12", "Caucasus & Middle East"),
    ("r13", "Central Asia"),
    ("r14", "South Asia West"),
    ("r15", "South Asia East"),
    ("r16", "Low Latitudes"),
    ("r17", "Southern Andes"),
    ("r18", "New Zealand"),
    ("r19", "Antarctic & Subantarctic"),
]

# ---------------------------------------------------------------------------
# Per-region slide content
# Each entry: (slide_A_bullets, slide_B_bullets, slide_C_bullets)
# ---------------------------------------------------------------------------
CONTENT = {
"r01": (
    # Slide A — time series
    [
        "One of the largest glacierized regions (~87,000 km²); strong negative mass balance trend since the 1980s driven by Gulf of Alaska maritime warming",
        "High inter-annual variability linked to ENSO and the Pacific Decadal Oscillation (PDO); El Niño years bring warmer, drier conditions that amplify mass loss",
        "Tidewater glaciers and large valley systems dominate; calving flux and surface melt both contribute to regional mass loss",
    ],
    # Slide B — feature importance (global bar)
    [
        "Ablation-season temperature (T abl mean, 29%) is the primary driver — summer warming controls melt of both maritime coastal glaciers and large mountain icefields",
        "Accumulation-season solar radiation (Rad acc, 14%) captures interannual cloud-cover variability linked to North Pacific atmospheric patterns (PDO, ENSO)",
        "Glacier ceiling elevation (Zmax, 7%) reflects hypsometric control — glaciers with high firn zones retain accumulation that offsets ablation-zone losses",
        "Climate accounts for 79% of attribution; without time encoding, all temporal variation is attributed to year-to-year climate forcing rather than a trend",
    ],
    # Slide C — beeswarm (local variability)
    [
        "Low-elevation coastal glaciers show the strongest ablation temperature attribution; high-elevation interior glaciers are more radiation-sensitive",
        "Radiation attribution is asymmetric: anomalously clear summers amplify melt beyond the temperature-only prediction, especially for large exposed icefields",
        "Zmax spread in the beeswarm reveals that high-accumulation-zone glaciers buffer mass loss — a topographic hedge against warming",
    ],
),
"r02": (
    [
        "~14,500 km² in Coast Mountains, Rockies, and Columbia Icefield; dramatic retreat since the 1980s including iconic glaciers (Athabasca, Illecillewaet)",
        "Precipitation-temperature balance defines a maritime (west coast) versus continental (Rockies) divide in glacier mass balance regime",
        "Strong mass loss acceleration in recent decades; Columbia Icefield glaciers particularly sensitive to summer temperature extremes",
    ],
    [
        "Accumulation-season solar radiation leads (Rad acc, 25%): Pacific storm variability controls winter cloud cover — clear vs cloudy winters shift both accumulation and energy balance simultaneously",
        "Ablation-season temperature (T abl, 19%): summer warmth drives melt on coastal maritime glaciers; Rockies glaciers are also vulnerable to extreme heat events",
        "Temperature variability (T abl var, 7%) captures heat-wave sensitivity — extreme melt events during anomalous summers disproportionately increase annual mass loss",
        "Median elevation (Zmed, 8%) encodes the maritime-continental divide: higher-elevation continental glaciers are less exposed to Pacific moisture and warm maritime air",
    ],
    [
        "Low-elevation coastal glaciers dominate the temperature attribution; high-elevation continental glaciers show more radiation sensitivity reflecting clearer continental skies",
        "Solar radiation attribution differentiates: cloudy maritime winters suppress Rad acc (positive for MB); clear continental summers amplify Rad abl (negative for MB)",
        "Precipitation (P acc) shows positive attribution, especially for coastal glaciers — high-snowfall winters offset summer melt, consistent with observed ELA dynamics",
    ],
),
"r03": (
    [
        "Large ice caps on Ellesmere and Devon Islands (~105,000 km²); some of the largest land ice outside Greenland and Antarctica",
        "Melt season is brief but intense; Arctic amplification means summer temperatures have risen faster here than the global average, driving accelerating mass loss",
        "High inter-annual variability tied to summer atmospheric circulation — blocking highs over the Arctic can trigger extreme melt seasons",
    ],
    [
        "Ablation-season temperature mean (T abl mean, 17%) AND variability (T abl var, 13%) are both in the top 2 — extreme warm summers during a brief Arctic melt window have a disproportionate impact on annual mass balance",
        "Combined radiation (Rad acc + Rad abl, ~19%): long Arctic summer days make solar radiation a co-driver of melt alongside temperature, especially on high-albedo plateau ice caps",
        "Arctic amplification sensitivity: small increases in summer temperature translate to large mass loss — the temperature signal is compressed into a short ablation window",
        "Geometry (23%) reflects ice cap plateau dynamics — large flat catchments respond collectively to temperature shifts, unlike valley glaciers",
    ],
    [
        "Temperature variability (T abl var) drives the widest spread in the beeswarm — melt event sensitivity varies strongly between northern high-elevation ice caps and lower southern margins",
        "Radiation drives mass loss preferentially for north-facing ice cap margins where temperature is marginal for melt; geometry buffers plateau interiors",
        "Local variability lower than expected for a large region — ice cap plateau glaciers respond more uniformly to climate than topographically complex valley systems",
    ],
),
"r04": (
    [
        "Baffin Island ice caps and valley glaciers (~40,000 km²); Barnes and Penny Ice Caps are among the last remnants of the Laurentide Ice Sheet",
        "Both ice caps are near or past tipping points — internal ice temperatures affect stability; ongoing rapid mass loss since the 1990s",
        "More maritime influence than R03 due to Hudson Strait; precipitation seasonality is higher and ablation seasons are warmer",
    ],
    [
        "Ablation temperature dominates at 35% — the highest single-feature concentration among Arctic Canadian regions; Baffin Island ice caps are acutely summer-temperature-sensitive with long, exposed low-elevation margins",
        "Temperature variability (T abl var, 15%) reflects vulnerability to anomalously warm summers — the 2012 extreme melt event is a canonical example of how one hot summer can drive a decade of normal mass loss",
        "Accumulation solar radiation (Rad acc, 15%) captures spring energy preconditioning — how much solar energy reaches the snowpack before the melt season affects how rapidly ablation begins",
        "Top 3 features alone explain 65% of variability — an unusually concentrated signal pointing to a single dominant climate control",
    ],
    [
        "Ice cap plateaux versus valley glacier margins diverge clearly in the beeswarm: plateaux respond uniformly to temperature; valley margins show stronger local geometry modulation",
        "Radiation attribution in accumulation season is positive for some glaciers (more solar energy → earlier melt onset → longer ablation season), reflecting preconditioning effects",
        "Geometry (Zmed, Zmax) introduces within-region variability: higher ice cap summits experience cooler temperatures and shorter melt seasons despite the same regional forcing",
    ],
),
"r05": (
    [
        "~130,000 km² of glaciers ringing the Greenland Ice Sheet — distinct from the GIS itself; includes outlet glaciers, ice caps, and valley glaciers",
        "Mass loss is partly driven by oceanic warming (fjord circulation bringing Atlantic water) but captured here through atmospheric forcing proxies",
        "Greenland Blocking events — anticyclonic conditions trapping warm air — drive extreme melt years; inter-annual variability is high",
    ],
    [
        "Near co-equal solar radiation and ablation temperature (Rad acc 20%, T abl 19%) — this dual-driver pattern reflects the transitional Arctic-maritime climate where both energy inputs shape melt independently",
        "Accumulation precipitation (P acc, 12%) is unusually prominent: peripheral Greenland glaciers are sensitive to snowfall anomalies that determine whether the winter mass gain offsets summer melt",
        "Temperature variability (T abl var, 11%): Greenland Blocking amplifies warm spells over peripheral glaciers; inter-annual variability in blocking frequency drives year-to-year extremes",
        "86% climate-driven — one of the most purely climate-forced regions globally, with limited topographic modulation",
    ],
    [
        "Beeswarm shows distinct east-west contrast: west Greenland (clearer skies) shows stronger radiation attribution; east Greenland (more cloud) shows stronger temperature attribution",
        "Precipitation attribution is consistently positive — more winter snowfall always improves mass balance — but magnitude varies with elevation",
        "Local temperature variability attribution identifies the glaciers most exposed to Greenland Blocking: low-latitude peripheral outlets in the southwest",
    ],
),
"r06": (
    [
        "~11,000 km² dominated by large ice caps (Vatnajökull, Langjökull, Hofsjökull); Iceland's glaciers have been losing mass since the late 1980s",
        "The North Atlantic Oscillation (NAO) is the dominant climate control: NAO+ winters deliver warmer temperatures AND more Atlantic precipitation simultaneously",
        "Geothermal heating under Vatnajökull adds a non-climatic mass loss component not captured by the model features",
    ],
    [
        "Ablation temperature leads (T abl, 18%) but accumulation precipitation is 2nd (P acc, 13%): Iceland's maritime mass balance is strongly accumulation-sensitive — NAO+ winters bring Atlantic moisture that can offset warmer temperatures",
        "Accumulation temperature (T acc, 10%) determines precipitation phase — mild Atlantic winters cause rain instead of snow at glacier margins, reducing accumulation even in wet years",
        "Multiple seasonal temperature variables (ablation, accumulation, variance) all in top 6: broad year-round sensitivity to North Atlantic climate variability driven by NAO and Atlantic Multidecadal Oscillation (AMO)",
        "Geometry contributes 25%, reflecting Iceland's complex ice cap drainage basin structures and diverse outlet glacier dynamics",
    ],
    [
        "Local beeswarm differentiates outlet glaciers (most temperature-sensitive) from high-elevation accumulation zones (most precipitation-sensitive)",
        "NAO co-variability visible in joint attribution of P acc and T acc: warm wet winters (NAO+) create opposing signals — more snow (positive) but rain risk at lower elevations (negative)",
        "Geothermal outlets (e.g. south Vatnajökull) may show anomalous patterns: lower temperature attribution because geothermal warming accelerates melt independently of atmospheric temperature",
    ],
),
"r07": (
    [
        "~34,000 km²; Svalbard is warming 3–4× the global average — one of the fastest-warming places on Earth — driven by Atlantification of the Arctic and sea ice loss feedbacks",
        "Mix of cold-based and polythermal glaciers; surge dynamics complicate long-term trends but the overall mass loss signal is strongly negative",
        "Sea ice decline creates a positive feedback: less sea ice → warmer surface air → more glacier melt",
    ],
    [
        "Ablation temperature (T abl mean) accounts for 41% of all attribution — the highest single-feature concentration of any region globally; Svalbard glaciers are in a melt-dominated regime where summer temperatures above freezing determine almost everything",
        "Temperature variability (T abl var, 13%): the frequency of extreme warm spells, amplified by reduced Arctic sea ice cover, captures the Atlantification signal — warmer Atlantic water drives both ocean and atmospheric temperature anomalies",
        "Median elevation (Zmed, 9%) modulates how much of each glacier's area falls below the equilibrium line altitude — the one topographic factor that differentiates melt exposure in a uniformly warm archipelago",
        "Accumulation precipitation (P acc, 9%): Svalbard receives Atlantic moisture; high-snowfall winters temporarily offset the strong melt trend",
    ],
    [
        "Beeswarm shows near-universal and consistent positive T abl attribution — Svalbard's climate forcing is remarkably spatially uniform, so within-region variability is lower than more topographically complex regions",
        "Zmed spread in the beeswarm is the main source of local variability: high-elevation glaciers retain more of their accumulation area; low-elevation tidewater glaciers lose mass from both calving and melt",
        "Marine-terminating glaciers may show weaker temperature attribution (calving losses not fully captured by atmospheric temperature) versus land-terminating glaciers",
    ],
),
"r08": (
    [
        "~2,700 km² with a strong climate gradient from maritime west-coast glaciers (Jostedalsbreen) to more continental interior glaciers further east",
        "NAO drives simultaneous temperature and precipitation anomalies — NAO+ winters are warmer AND wetter, creating opposing mass balance signals that partially cancel",
        "Mass loss accelerating since the 1990s; some Norwegian glaciers showed temporary advance in the early 2000s (NAO+ accumulation surplus) before resuming retreat",
    ],
    [
        "Near co-equal ablation temperature (T abl, 25%) and accumulation solar radiation (Rad acc, 24%) — a dual-driver regime reflecting the maritime fjord climate where cloud-cover variability and temperature anomalies are near-equally important",
        "Accumulation solar radiation modulates winter energy balance and snow metamorphism: clear winters accelerate spring snowpack ripening, extending the effective melt season — a mechanism specific to high-albedo maritime environments",
        "Multiple seasonal temperature variables (T acc mean, T acc var) also in top 8: NAO controls both ablation-season warmth and accumulation-season rain-versus-snow partitioning simultaneously",
        "The 25%/24% temperature-radiation near-tie is consistent with observed sensitivity of Norwegian glaciers to both atmospheric energy balance and temperature anomalies",
    ],
    [
        "West-coast maritime glaciers show strong radiation-temperature co-attribution; interior continental glaciers show stronger temperature dependence with weaker radiation signal",
        "Precipitation attribution is positive for west-coast glaciers (NAO+ → more Atlantic snowfall); near-zero or negative for continental glaciers where NAO controls temperature more than precipitation",
        "Local beeswarm should show bimodal structure corresponding to the maritime-continental divide — a clear geographic gradient in climate sensitivity",
    ],
),
"r09": (
    [
        "~51,500 km² scattered across Novaya Zemlya, Severnaya Zemlya, Franz Josef Land, and smaller Arctic islands — a highly heterogeneous archipelago setting",
        "Recent strong mass loss, particularly on Novaya Zemlya which hosts some of Russia's largest glaciers and ice caps; accelerating since ~2005",
        "Highly dispersed geography creates very different glaciological settings — from large Novaya Zemlya ice caps to small Franz Josef Land ice domes",
    ],
    [
        "The only region where geometry dominates over climate (51% vs 49%): extreme diversity of glacier types, sizes, and elevations across the archipelago disperses attribution across six geometric variables simultaneously",
        "Aspect (8%) is unusually prominent — north-facing vs south-facing orientation in these high-latitude archipelagos controls solar radiation exposure dramatically, creating strong within-region mass balance contrasts",
        "Ablation temperature remains the top single feature (21%): the shared large-scale Arctic warming still dominates individually, but geometric heterogeneity dilutes any single climate feature's fractional contribution",
        "log_Area (9%), Zmed (13%), and Zmin (7%) all prominent — glacier size, median elevation, and terminus position together determine which glaciers are most vulnerable to the common temperature forcing",
    ],
    [
        "Beeswarm shows the widest geometry attribution spread of any region: large Novaya Zemlya ice caps behave fundamentally differently from small Franz Josef Land ice domes in the same temperature regime",
        "Aspect attribution is directional: south-facing glaciers (positive Aspect value in northern hemisphere convention) show consistently stronger negative mass balance attribution — more insolation drives more melt at high latitudes",
        "Temperature attribution is compressed into a narrow range (less variable locally) while geometry attribution spans the full range — suggesting temperature sets the regional baseline but topography determines who loses most",
    ],
),
"r10": (
    [
        "~3,000 km² across Altai, Sayan, Stanovoy, and East Asian ranges — a dispersed continental glacier system with high seasonal temperature extremes",
        "Glaciers sit in a strongly continental climate with very clear skies; high summer solar radiation load and large temperature ranges characterise the region",
        "Debris cover and permafrost interactions on many glaciers create complex mass balance signals not directly captured by the model features",
    ],
    [
        "Accumulation-season solar radiation leads (Rad acc, 21%): the continental climate delivers high winter and spring insolation; solar energy controls snowpack energy balance during accumulation and determines spring melt onset timing",
        "Ablation temperature (T abl, 13%) and median elevation (Zmed, 12%) are nearly tied for 2nd: high continentality produces large temperature extremes, and glacier elevation determines the local melt season duration",
        "Lower elevation boundary (Zmin, 10%) is unusually prominent — glacier terminus position determines ablation zone extent; some debris-covered termini in the Altai effectively decouple terminus position from temperature forcing",
        "High geometry contribution (35%) reflects continental high-altitude settings where topographic position mediates very different local climate exposures across mountain ranges",
    ],
    [
        "Large local spread in Zmed and Zmin attributions — diverse settings from high Altai (termini at ~2,500m) to lower Sayan ranges (termini at ~1,800m) create very different melt exposures",
        "Solar radiation attribution is consistently negative (high insolation → enhanced melt and sublimation); Zmin attribution shows glaciers with higher termini lose less mass in the same radiation regime",
        "Debris-covered glacier termini may show weaker temperature and radiation attribution — insulation from debris suppresses melt response to the climate signal",
    ],
),
"r11": (
    [
        "~2,100 km² in the Alps, Pyrenees, and Apennines; Alpine glaciers have been retreating dramatically since the 1980s and are bellwethers of European climate change",
        "Multiple drivers: summer warming, declining winter snowfall, and increasing solar radiation (brightening) all contribute to observed mass loss",
        "Some Alpine glaciers may disappear entirely by 2100 under current trajectories; the loss of these glaciers will affect summer runoff for downstream populations",
    ],
    [
        "Solar radiation in both seasons (Rad acc 17% + Rad abl 10% = 27% combined) is the dominant energy input: under clear Alpine skies, direct solar radiation amplifies melt beyond what temperature alone predicts — a key distinction from higher-latitude regions",
        "Near-equal importance of ablation and accumulation season temperatures (T abl 14%, T acc 14%): broad seasonal sensitivity reflects the transitional continental-maritime Alpine climate — both cold winters and hot summers shape mass balance",
        "Temperature variability in accumulation season (T acc var, 9%) captures the frequency of mild winters with rain-on-snow events, which destroy the snowpack prematurely and extend the effective ablation season",
        "High climate dominance (83%) reflects the relatively homogeneous topographic character of Alpine glaciers — small, geometrically similar ice bodies are uniformly climate-forced",
    ],
    [
        "South-facing and low-elevation glaciers show the strongest combined radiation and temperature attribution — the hottest, brightest settings drive the largest mass loss",
        "Precipitation attribution is positive: years with heavy winter snowfall offset summer melt, consistent with observed ELA dynamics in the Alps where accumulation anomalies can outweigh temperature forcing in some years",
        "Accumulation temperature variability attribution reveals rain-on-snow event sensitivity — mild winter episodes that deposit rain instead of snow at glacier elevations leave a clear signal in the beeswarm",
    ],
),
"r12": (
    [
        "~1,400 km² in the Greater and Lesser Caucasus plus Turkish and Iranian highlands; glaciers sit at relatively low elevations for their latitude, making them acutely melt-sensitive",
        "Strong precipitation gradient: Atlantic and Black Sea moisture systems deliver winter snowfall to north-facing Greater Caucasus slopes; southern and eastern ranges are much drier",
        "Mass loss strong and consistent since the 1980s; some smaller glaciers in the Lesser Caucasus and Iranian highlands at risk of complete disappearance",
    ],
    [
        "Ablation temperature leads clearly (T abl mean, 25%): Caucasus glaciers sit at relatively low elevations (many below 3,500m); summer temperatures well above freezing drive prolonged and intense melt",
        "Median elevation (Zmed, 11%): the Greater Caucasus hosts glaciers spanning 2,000–5,600m; median elevation determines the length of the ablation season — the single most important topographic variable",
        "Accumulation precipitation (P acc, 9%): Winter Westerly Disturbances (WWDs) and Atlantic weather systems deliver the majority of snowfall to north-facing Caucasus slopes; variability in these systems drives accumulation anomalies",
        "Geometry contributes 28%, reflecting the steep terrain and diverse glacier types — from large valley glaciers on Mount Elbrus to small cirque glaciers in the Lesser Caucasus",
    ],
    [
        "North-facing Greater Caucasus glaciers (positive Aspect) show stronger precipitation attribution — they receive more winter snowfall from Atlantic moisture; south-facing glaciers are more temperature-dominated",
        "Elevation gradient is clear in the beeswarm: low-elevation glaciers show very strong positive T abl attribution (long melt seasons); high-elevation glaciers show more balanced temperature-radiation signals",
        "Precipitation attribution is positive and consistent for north-facing glaciers — Black Sea and Atlantic moisture input is a reliable accumulation driver, buffering temperature-driven mass loss",
    ],
),
"r13": (
    [
        "~70,000 km² spanning the Tian Shan, Pamir, Karakoram, and Hindu Kush — the 'Third Pole', a critical freshwater reservoir for millions in the Aral Sea, Indus, and Amu Darya basins",
        "Diverse mass balance regimes: summer-accumulation glaciers (Karakoram) coexist with winter-accumulation systems (western Tian Shan); the Karakoram Anomaly (stable or growing glaciers) complicates the regional signal",
        "Strong dependence on Westerly storm tracks for winter snowfall; Westerly variability linked to large-scale modes of variability (Arctic Oscillation, ENSO teleconnections)",
    ],
    [
        "Accumulation solar radiation leads (Rad acc, 23%): high-altitude continental setting with intense clear-sky insolation; winter solar energy controls snowpack metamorphism rates and spring melt onset — unique to high-altitude arid climates",
        "Ablation temperature (T abl, 15%): summer warming drives melt season length; the vast elevation range (2,000–7,000m+) means even modest warming substantially expands the ablation zone",
        "Westerly precipitation (P acc + P abl combined, 15%): Winter Westerly Disturbances deliver moisture to the Pamir and Tian Shan; inter-annual variability in Westerly intensity is a key accumulation driver — La Niña strengthens the Westerlies and increases Central Asian snowfall",
        "Elevation (Zmed, 11%) is the main geometry control: ELA position shifts dramatically with small temperature changes across the vast elevation range",
    ],
    [
        "Karakoram Anomaly visible in local patterns: some glaciers show unusual (positive or near-zero) temperature attribution in the expected ablation direction — reflecting the atypical cooling/thickening signal in some Karakoram valleys",
        "Summer-accumulation (Karakoram) versus winter-accumulation (Tian Shan) contrast: P abl attribution is positive for Karakoram glaciers (summer snowfall buffers melt) but not for western Tian Shan (summer rain events do not contribute to accumulation)",
        "Westerly precipitation attribution is spatially structured — glaciers facing Westerly storm tracks (northwest-facing, higher elevation) show stronger P acc attribution than rain-shadow easterly slopes",
    ],
),
"r14": (
    [
        "~33,000 km² in the western Himalaya (Indus basin), Karakoram, and Hindu Kush; dominated by Winter Westerly Disturbances (WWDs) for snowfall — the Indian Summer Monsoon penetration is limited here",
        "The Karakoram Anomaly is most pronounced in this region — some glaciers have been stable or thickening while global trends show mass loss",
        "Water tower of Pakistan: Indus river flow is critically dependent on glacier melt; changes here have direct consequences for downstream water security",
    ],
    [
        "Ablation temperature dominates at 36% — the most concentrated single driver outside of Svalbard; western Himalayan glaciers are acutely sensitive to summer warming because ablation is the dominant mass balance process",
        "Accumulation temperature (T acc, 16%): winter temperature determines whether WWD precipitation falls as snow or rain at glacier elevations — a critical phase transition; warmer winters with rain instead of snow dramatically reduce accumulation (Western Disturbance warming effect)",
        "Westerly influence encoded in both temperature and precipitation features: WWDs deliver the majority of high-elevation snowfall; the 51% combined temperature signal partly reflects that warm WWD events simultaneously deliver rain rather than snow",
        "Very low geometry contribution (14%): western Himalaya/Karakoram glaciers show relatively uniform climate sensitivity, suggesting topographic variability is secondary to large-scale climate forcing",
    ],
    [
        "Karakoram glaciers show anomalous temperature attribution patterns locally — some valleys experience cooling in summer due to circulation changes, creating positive T abl attribution (cooler → less melt) rather than the expected negative direction",
        "WWD precipitation attribution varies with northwest-facing exposure: glaciers on windward slopes of the Hindu Kush show strong positive P acc attribution; leeward glaciers receive less WWD moisture and show weaker precipitation signals",
        "The contrast between summer-melt-dominated Himalayan glaciers and dynamically active Karakoram glaciers (with surges) creates large local spread in the beeswarm",
    ],
),
"r15": (
    [
        "~14,700 km² in the eastern Himalaya and Hindu Kush (Brahmaputra and Ganges basins); uniquely dominated by the Indian Summer Monsoon (ISM) for both accumulation AND ablation",
        "A summer-accumulation regime: monsoon precipitation falls during the melt season; high-altitude snowfall during summer partially offsets melt, while low-altitude monsoon rain accelerates ablation",
        "Strong ISM inter-annual variability drives high year-to-year mass balance variability; ENSO modulates ISM intensity and thereby glacier mass balance",
    ],
    [
        "Ablation temperature leads (T abl, 31%) reflecting the dominant role of summer warming, but secondary features are more diversified than R14 — multiple energy inputs modulate the monsoon-season melt simultaneously",
        "Monsoon dual role: accumulation-season precipitation (P acc, 9%) AND ablation-season precipitation (P abl, 8%) are near-equally important — monsoon rain/snow events fall during the ablation season; summer snowfall above ~5,000m partially offsets melt while summer rain below that elevation accelerates it",
        "Solar radiation (Rad acc + Rad abl, 21% combined): monsoon cloud cover suppresses radiation during summer; interannual variability in monsoon cloud thickness modulates the energy available for melt, making radiation a secondary but significant driver",
        "In contrast to R14 (Westerly-driven), both precipitation features reflect the ISM: moisture arrives in summer and the rain/snow partition at glacier elevation determines whether precipitation adds or removes mass",
    ],
    [
        "Summer snowfall on high-elevation glaciers is visible in local patterns: P abl attribution is positive for glaciers above ~5,000m (monsoon snowfall buffers melt) and negative for lower-elevation glaciers (monsoon rain accelerates melt)",
        "ISM inter-annual variability creates wider local spread in precipitation attribution than in R14 — a strong monsoon year affects glaciers at all elevations simultaneously, though in opposite directions depending on elevation",
        "Temperature attribution is consistently dominant but more variable than in R14 — the noisy ISM signal introduces more climate uncertainty into the summer melt estimate",
    ],
),
"r16": (
    [
        "~3,000 km² of tropical glaciers in the tropical Andes (Peru, Bolivia, Ecuador, Colombia), East Africa (Kilimanjaro, Mt Kenya, Rwenzori), and Papua New Guinea",
        "These glaciers exist near the year-round freezing level — there is no distinct melt and accumulation season; melt is continuous and modulated by interannual ENSO-driven temperature and precipitation anomalies",
        "Among the most vulnerable glaciers globally; Kilimanjaro's ice fields have shrunk by ~80% since 1900; many tropical glaciers may disappear within decades",
    ],
    [
        "Accumulation temperature leads over ablation temperature (T acc 19% > T abl 15%) — unique among all 19 regions: at tropical glacier elevations near the year-round freezing point, temperature determines whether precipitation falls as rain or snow regardless of season, making accumulation-season temperature a critical phase-partition driver",
        "All 8 climate features appear in the top 8 with substantial contributions — the broadest climate distribution of any region, reflecting year-round climate sensitivity without a dominant season",
        "Solar radiation in both seasons (Rad acc + Rad abl, 25% combined): tropical glaciers receive intense, relatively consistent solar radiation; inter-annual variability in cloud cover (linked to ENSO sea surface temperatures) strongly modulates sublimation rates and energy available for melt",
        "Geometry contributes only 11%: these small, geometrically homogeneous high-elevation ice bodies are uniformly climate-controlled — topographic variability is secondary",
    ],
    [
        "Local patterns show elevation-dependent phase transition: glaciers at slightly lower elevations show T abl dominance (temperature frequently above freezing, direct melt control); higher-elevation glaciers show T acc dominance (temperature near 0°C determines rain vs snow)",
        "ENSO signature visible: El Niño years bring warmer, drier conditions to the tropical Andes (reduced P acc, increased T abl) — both attributions pull negative simultaneously, amplifying mass loss",
        "Precipitation attribution is positive for all glaciers — more precipitation always helps whether it falls as snow (accumulation) or the cooling effect of precipitation evaporation; the sign of T acc attribution reveals whether that precipitation phase is beneficial",
    ],
),
"r17": (
    [
        "~29,400 km²; dominated by the Northern and Southern Patagonian Ice Fields — the largest glaciers in the Southern Hemisphere outside Antarctica",
        "Extreme west-east precipitation gradient: Andean windward slopes receive >10 m/yr of precipitation from westerly storms; leeward (Argentine) slopes receive <1 m/yr",
        "Southern Annular Mode (SAM) controls westerly wind intensity and thereby cloud cover and precipitation delivery to the ice fields",
    ],
    [
        "Accumulation solar radiation leads (Rad acc, 21%): Patagonian cloud cover variability, driven by SAM and the intensity of the Southern Hemisphere westerly belt, is the dominant interannual mass balance signal — cloudy years (stronger westerlies) bring more precipitation and less radiation",
        "Ablation temperature (T abl, 17%): summer warming drives calving and surface melt at ice field margins; lake-calving glaciers (Grey, Perito Moreno, etc.) are particularly sensitive to summer atmospheric temperature",
        "Median elevation (Zmed, 12%) is the key geometry control — orographic position (windward vs leeward of the Andean divide) is encoded by elevation: low-elevation western outlets are precipitation-dominated; higher-elevation eastern outlets are temperature-dominated",
        "Ablation-season precipitation (P abl, 9%): Southern Andes receive precipitation year-round from westerly systems; summer rain events affect glacier albedo and surface energy balance — positive SAM drives year-round storm track shifts",
    ],
    [
        "Strong east-west contrast in local attribution: windward (western) glaciers show strong radiation-precipitation co-attribution; leeward (eastern) glaciers show stronger temperature sensitivity",
        "SAM modulation is visible: positive SAM → stronger westerlies → more cloud cover on windward side → less solar radiation → less melt; captured in the negative relationship between Rad acc attribution and SAM index years",
        "Large calving glaciers (Upsala, Perito Moreno) may show weaker temperature attribution than land-terminating glaciers — calving dynamics add a non-atmospheric mass loss component",
    ],
),
"r18": (
    [
        "~1,160 km² in the Southern Alps of New Zealand's South Island; among the most maritime glaciers globally, receiving extreme precipitation from the Roaring Forties",
        "Franz Josef and Fox Glaciers are unusually responsive to climate — their rapid advance-retreat cycles on decadal timescales are driven by ENSO and SAM-modulated precipitation and temperature",
        "ENSO is the dominant climate driver: La Niña → more southwesterly flow → more precipitation and positive MB; El Niño → warmer, drier → negative MB",
    ],
    [
        "Temperature and precipitation are near co-equal primary drivers, reflecting New Zealand's extreme maritime regime where both energy (temperature) and mass (precipitation) forcing are large and variable",
        "Combined precipitation (P acc + P abl, ~24%) is the highest precipitation fraction of any region outside Iceland, reflecting the extreme precipitation regime from the Roaring Forties — interannual variability in storm track position is as important as temperature",
        "Both accumulation AND ablation season precipitation are prominent: New Zealand glaciers receive significant year-round precipitation; ablation-season precipitation phase (rain vs snow above ~1,500m) is temperature-dependent and drives strong co-variability between T and P features",
        "ENSO and SAM act simultaneously through both temperature and precipitation channels — ENSO La Niña both cools AND increases precipitation; this co-variability is captured implicitly through the correlated T-P feature attributions",
    ],
    [
        "Low-elevation glaciers (Franz Josef, Fox termini) are most temperature-sensitive — melt is continuous, and temperature controls the rate directly; high-elevation glaciers show more precipitation sensitivity",
        "Positive precipitation attribution is universal — more snowfall always benefits mass balance; the interaction with T acc attribution identifies whether that precipitation falls as snow (cold) or rain (warm) at glacier elevation",
        "The wide local spread in both T and P attributions reflects the strong elevation gradient in New Zealand's compact Southern Alps — a few hundred metres of elevation shifts the climate regime from maritime-dominated to radiation-controlled",
    ],
),
"r19": (
    [
        "~132,000 km²; covers South Georgia, Kerguelen, Heard Island, South Sandwich Islands, Falkland Islands glaciers, and Antarctic Peninsula glaciers — an extremely diverse set of islands spanning 45°S–65°S",
        "Southern Ocean storm tracks and the Southern Annular Mode (SAM) control climate across these islands; each island sits in a different position relative to the circumpolar westerly belt",
        "Mass loss rates vary enormously across the region — South Georgia is losing mass rapidly; some smaller Subantarctic islands show more stable conditions depending on their SAM exposure",
    ],
    [
        "Second-highest geometry contribution (36%) after Russian Arctic: extreme structural diversity from large ice sheets (South Georgia, Heard Island) to small cirque glaciers on Kerguelen and the South Sandwich Islands demands topographic context to interpret climate sensitivity",
        "Glacier size (log_Area, 9%) and slope (7%) are unusual top-8 features — large ice bodies on South Georgia behave as ice sheets with internal dynamics; small cirque glaciers on oceanic islands behave as direct climate recorders",
        "Ablation temperature leads (T abl, 22%) with accumulation temperature (T acc, 14%) also prominent — subpolar marine climate means atmospheric temperature is important year-round, with no prolonged freeze period on lower islands",
        "Subantarctic precipitation (P acc, 8%) reflects Southern Ocean storm track variability linked to SAM and ENSO teleconnections — SAM modulates westerly intensity and thereby precipitation delivery to each island",
    ],
    [
        "Enormous local variability in geometry attribution: South Georgia's large ice masses show strong log_Area and slope attribution; small Kerguelen island glaciers show much more climate-dominated, geometry-independent signals",
        "SAM influence visible through co-variability of radiation and precipitation attributions: positive SAM shifts storm tracks poleward, increasing precipitation on some islands (South Georgia) while reducing it on others (Kerguelen, at lower latitude) — creating a structured geographic pattern in P acc attribution",
        "Antarctic Peninsula glaciers in the northern extent of this region may show stronger temperature attribution from the rapid warming on the peninsula since the 1980s — distinct from the more stable subpolar Subantarctic islands",
    ],
),
}

# ---------------------------------------------------------------------------
# PPT building helpers
# ---------------------------------------------------------------------------

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs


def _blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(blank_layout)


def _add_rect(slide, l, t, w, h, fill_rgb: RGBColor, alpha=None):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape


def _add_text_box(slide, l, t, w, h, text: str, font_size: int,
                  bold=False, color=WHITE, align=PP_ALIGN.LEFT,
                  wrap=True, font_name="Calibri"):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.name = font_name
    return txb


def _add_bullet_box(slide, l, t, w, h, bullets: list[str],
                    font_size=10.5, color=DARK, header: str | None = None):
    from pptx.oxml.ns import qn
    from lxml import etree

    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True

    first = True
    if header:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = header
        run.font.bold = True
        run.font.size = Pt(font_size + 0.5)
        run.font.color.rgb = NAVY
        run.font.name = "Calibri"
        p.space_after = Pt(4)

    for bullet in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(3)
        p.space_after  = Pt(3)
        # bullet character
        run0 = p.add_run()
        run0.text = "▸  "
        run0.font.size = Pt(font_size - 1)
        run0.font.color.rgb = TEAL
        run0.font.name = "Calibri"
        run0.font.bold = True
        run = p.add_run()
        run.text = bullet
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return txb


def _add_title_bar(slide, title_text: str, rkey: str, region_name: str,
                   slide_type: str):
    """Dark navy title bar spanning full width."""
    _add_rect(slide, Inches(0), Inches(0), SW, Inches(1.0), NAVY)

    # Region tag (left, small)
    _add_text_box(slide,
                  Inches(0.25), Inches(0.06), Inches(3), Inches(0.35),
                  f"{rkey.upper()}  ·  {region_name}",
                  font_size=9, bold=False, color=SILVER)

    # Slide type tag (right)
    _add_text_box(slide,
                  Inches(11.5), Inches(0.06), Inches(1.6), Inches(0.35),
                  slide_type,
                  font_size=9, bold=True, color=GROUP, align=PP_ALIGN.RIGHT)

    # Main title
    _add_text_box(slide,
                  Inches(0.25), Inches(0.38), Inches(12.8), Inches(0.55),
                  title_text,
                  font_size=16, bold=True, color=WHITE)


def _add_footer(slide, text="Glacier Mass Balance · BayesNF · No Time Encoding"):
    _add_rect(slide, Inches(0), Inches(7.1), SW, Inches(0.4), NAVY)
    _add_text_box(slide,
                  Inches(0.3), Inches(7.12), Inches(12.5), Inches(0.32),
                  text, font_size=8, color=SILVER)


def add_title_slide(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _add_rect(slide, Inches(0), Inches(0), SW, SH, NAVY)
    _add_rect(slide, Inches(0), Inches(3.2), SW, Inches(0.06), TEAL)

    _add_text_box(slide,
                  Inches(1.0), Inches(1.2), Inches(11.3), Inches(1.4),
                  "Regional Glacier Mass Balance",
                  font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_text_box(slide,
                  Inches(1.0), Inches(2.6), Inches(11.3), Inches(0.7),
                  "Feature Importance & Physical Interpretation — No Time Encoding",
                  font_size=18, bold=False, color=TEAL, align=PP_ALIGN.CENTER)
    _add_text_box(slide,
                  Inches(1.0), Inches(3.5), Inches(11.3), Inches(0.5),
                  "All 19 RGI Regions  ·  BayesNF (Flax/JAX)  ·  Ensemble Top-5 by GLaMBIE RMSE",
                  font_size=13, bold=False, color=SILVER, align=PP_ALIGN.CENTER)
    _add_text_box(slide,
                  Inches(1.0), Inches(6.7), Inches(11.3), Inches(0.4),
                  "Ritu Anilkumar  ·  Jonathan Bamber  ·  Fabien Maussion  ·  Michael Zemp",
                  font_size=11, bold=False, color=SILVER, align=PP_ALIGN.CENTER)


def add_region_slide_A(prs, rkey, region_name, img_path, bullets):
    slide = _blank_slide(prs)
    _add_rect(slide, Inches(0), Inches(0), SW, SH, OFFWHITE)
    _add_title_bar(slide, f"{region_name} — Regional Mass Balance Time Series",
                   rkey, region_name, "Slide A")
    slide.shapes.add_picture(str(img_path), IMG_L, IMG_T, IMG_W, IMG_H)
    _add_bullet_box(slide, BUL_L, BUL_T, BUL_W, BUL_H, bullets,
                    header="Regional context")
    _add_footer(slide)


def add_region_slide_B(prs, rkey, region_name, img_path, bullets):
    slide = _blank_slide(prs)
    _add_rect(slide, Inches(0), Inches(0), SW, SH, OFFWHITE)
    _add_title_bar(slide, f"{region_name} — Feature Importance (Global)",
                   rkey, region_name, "Slide B")
    slide.shapes.add_picture(str(img_path), IMG_L, IMG_T, IMG_W, IMG_H)
    _add_bullet_box(slide, BUL_L, BUL_T, BUL_W, BUL_H, bullets,
                    header="Feature importance")
    _add_footer(slide)


def add_region_slide_C(prs, rkey, region_name, img_path, bullets):
    slide = _blank_slide(prs)
    _add_rect(slide, Inches(0), Inches(0), SW, SH, OFFWHITE)
    _add_title_bar(slide, f"{region_name} — Feature Importance (Local Variability)",
                   rkey, region_name, "Slide C")
    slide.shapes.add_picture(str(img_path), IMG_L, IMG_T, IMG_W, IMG_H)
    _add_bullet_box(slide, BUL_L, BUL_T, BUL_W, BUL_H, bullets,
                    header="Local variability & physics")
    _add_footer(slide)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--ensemble_root", default="outputs/egu_fin/ensemble_te_r2")
    parser.add_argument("--explain_root",  default="outputs/egu_fin/explain_te")
    parser.add_argument("--group",         default="no_time_encoding")
    parser.add_argument("--output",        default="outputs/egu_fin/regional_no_te.pptx")
    args = parser.parse_args()

    ensemble_root = Path(args.ensemble_root)
    explain_root  = Path(args.explain_root)
    group         = args.group
    output        = Path(args.output)
    output.parent.mkdir(parents=True, exist_ok=True)

    prs = new_prs()
    add_title_slide(prs)

    for rkey, region_name in REGIONS:
        ts_img  = ensemble_root / rkey / group / "ensemble_regional_gt.png"
        bar_img = explain_root  / rkey / group / f"importance_bar_finetune_ensemble1_{group}_masked.png"
        bee_img = explain_root  / rkey / group / f"beeswarm_finetune_ensemble1_{group}_masked.png"

        missing = [p for p in [ts_img, bar_img, bee_img] if not p.exists()]
        if missing:
            print(f"  [{rkey}] MISSING: {[str(p) for p in missing]}")
            continue

        bullets_A, bullets_B, bullets_C = CONTENT[rkey]

        print(f"  Adding {rkey} — {region_name}")
        add_region_slide_A(prs, rkey, region_name, ts_img,  bullets_A)
        add_region_slide_B(prs, rkey, region_name, bar_img, bullets_B)
        add_region_slide_C(prs, rkey, region_name, bee_img, bullets_C)

    prs.save(str(output))
    print(f"\nSaved: {output}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
